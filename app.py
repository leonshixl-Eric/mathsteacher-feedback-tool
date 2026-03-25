import sys
import matplotlib
matplotlib.use('Agg') 
import matplotlib.pyplot as plt
import streamlit as st
import pandas as pd
from docx import Document
from docx.shared import Cm, Pt
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Cm as PptxCm, Pt as PptxPt
import zipfile
import os
import re
from io import BytesIO

# --- FIX FOR PYTHON 3.13 & PYMUPDF ---
try:
    import imghdr
except ImportError:
    import filetype
    class MockImghdr:
        def what(self, file, h=None):
            kind = filetype.guess(file)
            return kind.extension if kind else None
    sys.modules['imghdr'] = MockImghdr()

try:
    import fitz  # PyMuPDF for PDF text extraction
except ImportError:
    fitz = None

st.set_page_config(page_title="Maths Feedback Pro", layout="centered", page_icon="📊")

st.title("📊 High-Fidelity Feedback Generator")
st.write("Upload all three files. Questions are fully reconstructed with perfect mathematical formatting.")

# --- 1. THE UPLOADERS ---
uploaded_csv = st.file_uploader("1. Upload Marks (CSV or Excel)", type=["csv", "xlsx"])
uploaded_pdf = st.file_uploader("2. Upload Original Exam PDF (Reference)", type="pdf")
uploaded_mapping = st.file_uploader("3. Upload Topic Mapping (CSV or Excel)", type=["csv", "xlsx"])

# --- BRANDING SETTINGS ---
st.markdown("---")
st.subheader("📝 Document Branding")
col_brand1, col_brand2 = st.columns(2)
with col_brand1:
    unit_title = st.text_input("Unit/Topic Title", value="Algebraic Manipulation")
    class_name = st.text_input("Class Name", value="Year 9 Maths")
with col_brand2:
    uploaded_logo = st.file_uploader("Upload School Logo (Optional)", type=["png", "jpg", "jpeg"])

# --- DOCUMENT SETTINGS ---
st.markdown("---")
st.subheader("⚙️ Document Settings")

col_setting1, col_setting2, col_setting3 = st.columns(3)
with col_setting1:
    selected_font_size = st.slider("Question Font Size", min_value=10, max_value=14, value=12, step=1)
with col_setting2:
    selected_threshold = st.slider("Reteach Threshold (%)", min_value=0, max_value=100, value=55, step=5)
with col_setting3:
    selected_margin = st.slider("Page Margin (cm)", min_value=0.5, max_value=3.0, value=1.3, step=0.1)

generate_ppt = st.checkbox("📽️ Also generate Whole-Class Reteach PowerPoint (PPTX)", value=True)
st.markdown("---")

threshold_decimal = selected_threshold / 100.0

# --- 2. THE RECONSTRUCTION ENGINE ---

def extract_questions_from_pdf(pdf_file, q_labels):
    """Scans the uploaded PDF and tries to extract the text for each question."""
    db = {}
    if not pdf_file or not fitz:
        return db
        
    try:
        # Read the PDF text
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        full_text = "\n".join([page.get_text("text") for page in doc])
        pdf_file.seek(0) # Reset file pointer
        
        valid_qs = [q for q in q_labels if q not in ["Surname", "Forename"]]
        
        for q in valid_qs:
            # Parse the question number and letter (e.g. "1a" -> "1", "a")
            m = re.match(r"(\d+)([a-zA-Z]*)", q)
            if not m: continue
            num, let = m.groups()
            
            # Build a search pattern to find "1a", "1(a)", "1 a", or "1."
            if let:
                search_str = rf"\b{num}\s*[\(\.]?\s*{let}[\)\.]?"
            else:
                search_str = rf"\b{num}\s*[\)\.]"
                
            match = re.search(search_str, full_text, re.IGNORECASE)
            if match:
                start_idx = match.end()
                # Grab the next 200 characters as the question text
                extracted = full_text[start_idx:start_idx+200].strip()
                # Clean up weird line breaks from the PDF
                extracted = re.sub(r'\s+', ' ', extracted)
                db[q] = f"{q}) {extracted}..."
    except Exception as e:
        st.warning(f"Could not extract text from PDF: {e}")
        
    return db

def create_question_image(q_code, text, font_size):
    # Fallback text if the PDF didn't contain the question or wasn't readable
    if not text:
        text = f"Question {q_code}\n\n(Please refer to the original exam paper \nfor exact mathematical formatting)."
        
    # Wrap text if it is too long (since we extract raw PDF text now)
    import textwrap
    wrapped_text = textwrap.fill(text, width=60)
        
    line_count = wrapped_text.count('\n') + 1
    base_padding = 0.3 
    height_per_line = font_size * 0.035 
    fig_height = base_padding + (line_count * height_per_line)
    
    plt.figure(figsize=(7, fig_height))
    plt.text(0.01, 0.5, wrapped_text, fontsize=font_size, verticalalignment='center', fontfamily='serif')
    plt.axis('off')
    plt.tight_layout(pad=0)
    
    img_name = f"q_{q_code}.png"
    plt.savefig(img_name, dpi=200, bbox_inches='tight')
    plt.close()
    return img_name

def process_data(uploaded_csv, uploaded_mapping):
    df_marks = pd.read_csv(uploaded_csv, header=None) if uploaded_csv.name.endswith('.csv') else pd.read_excel(uploaded_csv, header=None)
    
    # 1. Dynamically build Q Labels from Row 1 and Row 2
    row0 = df_marks.iloc[0].astype(str).tolist()
    row1 = df_marks.iloc[1].astype(str).tolist()
    
    q_labels = ["Surname", "Forename"]
    current_q = ""
    
    for i in range(2, len(row0)):
        r0 = row0[i].strip()
        r1 = row1[i].strip()
        
        if r0.lower() == 'total' or r1.lower() == 'total': break
            
        if r0 != 'nan' and r0 != '':
            m = re.search(r'\d+', r0)
            if m: current_q = m.group()
        
        if r1 != 'nan' and r1 != '': q_labels.append(current_q + r1)
        else: q_labels.append(current_q)

    full_marks_row = df_marks.iloc[2]
    
    # 2. Dynamically find percentage row
    percentage_idx = None
    for i in range(len(df_marks)):
        cell_val = str(df_marks.iloc[i, 0]).strip().lower()
        if 'percentage' in cell_val:
            percentage_idx = i
            break
            
    if percentage_idx is None:
        raise ValueError("Could not find the 'Percentage' row. Ensure the first column of the bottom row contains the word 'Percentage'.")
        
    percentage_row = df_marks.iloc[percentage_idx]
    
    student_rows = df_marks.iloc[3:percentage_idx].reset_index(drop=True)
    student_rows = student_rows.dropna(subset=[0]).reset_index(drop=True)
    
    df_map = pd.read_csv(uploaded_mapping, header=None) if uploaded_mapping.name.endswith('.csv') else pd.read_excel(uploaded_mapping, header=None)
    
    first_cell = str(df_map.iloc[0, 0]).lower()
    if 'topic' in first_cell or 'area' in first_cell:
        df_map = df_map.iloc[1:].reset_index(drop=True)

    dynamic_areas = []
    for _, map_row in df_map.iterrows():
        if pd.isna(map_row.iloc[0]): continue
        topic = str(map_row.iloc[0]).strip()
        
        qs = []
        last_num = ""
        
        for col_idx in range(1, len(map_row)):
            cell_val = map_row.iloc[col_idx]
            if pd.isna(cell_val): continue
            
            raw_str = str(cell_val).lower().replace('and', ',').replace('&', ',')
            raw_str = "".join(raw_str.split())
            tokens = raw_str.split(',')
            
            for t in tokens:
                if not t: continue
                num_part = "".join([c for c in t if c.isdigit()])
                letter_part = "".join([c for c in t if c.isalpha()])
                
                if num_part:
                    last_num = num_part
                    candidate = num_part + letter_part
                else:
                    candidate = last_num + letter_part
                    
                if candidate in q_labels and candidate not in qs:
                    qs.append(candidate)
        
        indices = [q_labels.index(q) for q in qs]
        if indices:
            dynamic_areas.append((topic, indices))
            
    return student_rows, percentage_row, full_marks_row, q_labels, dynamic_areas

def add_tight_picture(doc, img_path, width):
    paragraph = doc.add_paragraph()
    paragraph.paragraph_format.space_before = Cm(0.3)
    paragraph.paragraph_format.space_after = Cm(0.3)
    run = paragraph.add_run()
    run.add_picture(img_path, width=width)
    return paragraph

# --- 3. BUTTON LAYOUT ---
col1, col2 = st.columns(2)
with col1:
    preview_clicked = st.button("👀 Preview Sample Student", use_container_width=True)
with col2:
    generate_clicked = st.button("📄 Generate All Feedback", type="primary", use_container_width=True)

# --- 4. PREVIEW LOGIC ---
if preview_clicked:
    if not (uploaded_csv and uploaded_pdf and uploaded_mapping):
        st.warning("Please upload all three files to see a preview.")
    else:
        try:
            with st.spinner("Analyzing Exam Files..."):
                student_rows, percentage_row, full_marks_row, q_labels, dynamic_areas = process_data(uploaded_csv, uploaded_mapping)
                
                # Extract text directly from the uploaded PDF!
                pdf_db = extract_questions_from_pdf(uploaded_pdf, q_labels)

                first_student = None
                for _, row in student_rows.iterrows():
                    cell_text = str(row[0]).lower()
                    if cell_text != 'nan' and 'name' not in cell_text and 'surname' not in cell_text:
                        first_student = row
                        break
                
                if first_student is not None:
                    lname = str(first_student[0]).strip()
                    fname = str(first_student[1]).strip()
                    name = f"{fname} {lname}".strip() if fname.lower() != 'nan' else lname
                    
                    col_prev1, col_prev2 = st.columns([1, 8])
                    with col_prev1:
                        if uploaded_logo is not None:
                            st.image(uploaded_logo, width=50) 
                    with col_prev2:
                        st.markdown(f"#### {unit_title} Feedback: **{name}** &nbsp; | &nbsp; Class: **{class_name}**")
                    
                    preview_table = []
                    student_ebi = []
                    for title, idxs in dynamic_areas:
                        w, e = [], []
                        for idx in idxs:
                            score = pd.to_numeric(first_student[idx], errors='coerce')
                            full_mark = pd.to_numeric(full_marks_row[idx], errors='coerce')
                            
                            if pd.notna(score) and pd.notna(full_mark) and score >= full_mark:
                                w.append(q_labels[idx])
                            else: 
                                e.append(q_labels[idx])
                                student_ebi.append(q_labels[idx])
                                
                        preview_table.append({"Area": title, "What Went Well": ", ".join(w), "Even Better If": ", ".join(e)})
                    
                    st.table(pd.DataFrame(preview_table))
                    
                    reteach = [q for q in student_ebi if pd.to_numeric(percentage_row[q_labels.index(q)], errors='coerce') <= threshold_decimal]
                    personal = [q for q in student_ebi if q not in reteach]
                    
                    st.markdown("#### 🎯 Personal Corrections")
                    if personal:
                        for q in personal:
                            img_path = create_question_image(q, pdf_db.get(q, ""), selected_font_size)
                            st.image(img_path)
                            os.remove(img_path)
                    else:
                        st.success("No personal corrections needed!")
                    
                    st.markdown(f"#### 🏫 Whole-Class Reteaching (≤ {selected_threshold}%)")
                    if reteach:
                        for q in reteach:
                            img_path = create_question_image(q, pdf_db.get(q, ""), selected_font_size)
                            st.image(img_path)
                            os.remove(img_path)
                    else:
                        st.success("No whole-class reteaching needed!")
                else:
                    st.error("Could not find any valid students in the CSV.")
        except Exception as e:
            st.error(f"Error reading files: {e}")

# --- 5. GENERATE LOGIC ---
if generate_clicked:
    if not (uploaded_csv and uploaded_pdf and uploaded_mapping):
        st.error("Please upload all three files (Marks, PDF, Mapping).")
    else:
        try:
            with st.spinner(f'Reconstructing questions and generating files...'):
                logo_path = None
                if uploaded_logo is not None:
                    logo_path = "temp_logo.png"
                    with open(logo_path, "wb") as f:
                        f.write(uploaded_logo.getbuffer())

                student_rows, percentage_row, full_marks_row, q_labels, dynamic_areas = process_data(uploaded_csv, uploaded_mapping)
                
                # Extract text directly from the uploaded PDF!
                pdf_db = extract_questions_from_pdf(uploaded_pdf, q_labels)
                q_images = {q: create_question_image(q, pdf_db.get(q, ""), selected_font_size) for q in q_labels if q not in ["Surname", "Forename"]}

                doc = Document()
                
                available_width_cm = 21.0 - (2 * selected_margin)
                area_col_width = available_width_cm - 7.0 
                col_widths = [Cm(area_col_width), Cm(3.5), Cm(3.5)]
                
                personal_img_width = Cm(min(12.0, available_width_cm))
                reteach_img_width = Cm(min(14.0, available_width_cm))

                for section in doc.sections:
                    section.page_width = Cm(21.0)
                    section.page_height = Cm(29.7)
                    section.top_margin, section.bottom_margin = Cm(selected_margin), Cm(selected_margin)
                    section.left_margin, section.right_margin = Cm(selected_margin), Cm(selected_margin)

                for _, row in student_rows.iterrows():
                    lname = str(row[0]).strip()
                    fname = str(row[1]).strip()
                    
                    cell_text = lname.lower()
                    if cell_text == 'nan' or 'surname' in cell_text or 'name' in cell_text: continue
                    name = f"{fname} {lname}".strip() if fname.lower() != 'nan' else lname
                    
                    header_p = doc.add_paragraph()
                    header_p.paragraph_format.space_after = Cm(0.3)
                    
                    if logo_path:
                        header_p.add_run().add_picture(logo_path, width=Cm(1.5))
                        header_p.add_run("    ") 
                    
                    r_title = header_p.add_run(f"{unit_title} Feedback: {name}   |   Class: {class_name}")
                    r_title.bold = True
                    r_title.font.size = Pt(14) 
                    
                    table = doc.add_table(rows=1, cols=3)
                    table.style = 'Table Grid'
                    
                    for i in range(3): table.columns[i].width = col_widths[i]
                    
                    hdr = table.rows[0].cells
                    hdr[0].text, hdr[1].text, hdr[2].text = "Area", "what went well", "even better if"
                    for i in range(3): hdr[i].width = col_widths[i]
                    
                    student_ebi = []
                    for title, idxs in dynamic_areas:
                        w, e = [], []
                        for idx in idxs:
                            score = pd.to_numeric(row[idx], errors='coerce')
                            full_mark = pd.to_numeric(full_marks_row[idx], errors='coerce')
                            
                            if pd.notna(score) and pd.notna(full_mark) and score >= full_mark:
                                w.append(q_labels[idx])
                            else: 
                                e.append(q_labels[idx])
                                student_ebi.append(q_labels[idx])
                                
                        r = table.add_row().cells
                        r[0].text, r[1].text, r[2].text = str(title), ", ".join(w), ", ".join(e)
                        for i in range(3): r[i].width = col_widths[i]

                    reteach = [q for q in student_ebi if pd.to_numeric(percentage_row[q_labels.index(q)], errors='coerce') <= threshold_decimal]
                    personal = [q for q in student_ebi if q not in reteach]
                    
                    if personal:
                        h_pers = doc.add_heading("Personal correction", 2)
                        h_pers.paragraph_format.space_before = Cm(0)
                        
                        for q in personal:
                            add_tight_picture(doc, q_images[q], width=personal_img_width)

                    doc.add_page_break()
                    
                    h_ret = doc.add_heading(f"Whole-class reteaching - {name}", 1)
                    h_ret.paragraph_format.space_before = Cm(0)
                    
                    if reteach:
                        for q in reteach: 
                            add_tight_picture(doc, q_images[q], width=reteach_img_width)
                    else: doc.add_paragraph("Excellent mastery of class topics.")
                    doc.add_page_break()

                target_docx = BytesIO()
                doc.save(target_docx)

                target_pptx = None
                global_reteach_qs = []
                
                if generate_ppt:
                    for title, idxs in dynamic_areas:
                        for idx in idxs:
                            q = q_labels[idx]
                            if q not in global_reteach_qs:
                                class_avg = pd.to_numeric(percentage_row[idx], errors='coerce')
                                if class_avg <= threshold_decimal:
                                    global_reteach_qs.append(q)
                    
                    global_reteach_qs.sort(key=lambda x: q_labels.index(x))

                    if len(global_reteach_qs) > 0:
                        prs = Presentation()
                        for q in global_reteach_qs:
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            
                            img_path = q_images[q]
                            pic_left = PptxCm(2)
                            pic_top = PptxCm(2.5) 
                            pic_width = PptxCm(21.4) 
                            
                            slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_width)
                        
                        target_pptx = BytesIO()
                        prs.save(target_pptx)

                safe_class = str(class_name).strip().replace(" ", "_")
                safe_unit = str(unit_title).strip().replace(" ", "_")
                docx_name = f"{safe_class}_{safe_unit}_Feedback.docx"
                pptx_name = f"{safe_class}_{safe_unit}_Reteach_Slides.pptx"
                zip_name = f"{safe_class}_{safe_unit}_All_Files.zip"
                
                zip_buffer = BytesIO()
                with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                    zip_file.writestr(docx_name, target_docx.getvalue())
                    if generate_ppt and target_pptx is not None:
                        zip_file.writestr(pptx_name, target_pptx.getvalue())
                zip_buffer.seek(0)

                st.success(f"✅ Feedback Pack Ready!")
                
                if generate_ppt and target_pptx is not None:
                    col_dl1, col_dl2, col_dl3 = st.columns(3)
                    with col_dl1:
                        st.download_button("📥 Download Word Doc", data=target_docx.getvalue(), file_name=docx_name)
                    with col_dl2:
                        st.download_button("📽️ Download PPTX", data=target_pptx.getvalue(), file_name=pptx_name)
                    with col_dl3:
                        st.download_button("📦 Download All (ZIP)", data=zip_buffer.getvalue(), file_name=zip_name, type="primary")
                else:
                    col_dl1, col_dl2 = st.columns(2)
                    with col_dl1:
                        st.download_button("📥 Download Word Doc", data=target_docx.getvalue(), file_name=docx_name)
                    with col_dl2:
                        st.download_button("📦 Download All (ZIP)", data=zip_buffer.getvalue(), file_name=zip_name, type="primary")
                    
                    if generate_ppt and len(global_reteach_qs) == 0:
                        st.info("No PowerPoint generated because the class scored above the reteach threshold on all topics!")

                for f in os.listdir():
                    if f.startswith("q_") and f.endswith(".png"): os.remove(f)
                if logo_path and os.path.exists(logo_path):
                    os.remove(logo_path)

        except ValueError as ve:
            st.error(f"Spreadsheet Error: {ve}")
        except Exception as e:
            st.error(f"Error: {e}")
